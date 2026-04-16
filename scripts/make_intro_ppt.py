"""
Police Frameworks — 소개 슬라이드 생성기 (60대 경찰관 대상)

구조 (27장):
    1        표지
    2        차례
    3-17     사례 5건 × 3장 (상황 / 분석 / 실행 방안)
    18       공구함 비유
    19       12개 공구 4서랍
    20       서랍 1 상세
    21       서랍 2·3 상세
    22       서랍 4 상세
    23       /peel 안내데스크
    24       Claude Skills 어디에 깔리나
    25       Claude에게 설치 부탁하기
    26       첫 대화 예시
    27       부탁과 연락처

작성자: 최희철 (경주경찰서 경찰발전협의회 회원, 소프트웨어 개발자)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ──────────────────────────────────────────────────────────────
# 색상
# ──────────────────────────────────────────────────────────────

NAVY = RGBColor(0x1A, 0x36, 0x5D)
NAVY_DEEP = RGBColor(0x0F, 0x23, 0x3F)
NAVY_LIGHT = RGBColor(0x2C, 0x51, 0x82)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x0F, 0x76, 0x42)
GOLD = RGBColor(0xD6, 0x9E, 0x2E)
GRAY_DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY_MID = RGBColor(0x4A, 0x55, 0x68)
GRAY_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_PALE = RGBColor(0xF1, 0xF5, 0xF9)
BG = RGBColor(0xFA, 0xFC, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW_HL = RGBColor(0xFE, 0xF3, 0xC7)

FONT = "맑은 고딕"
TOTAL_SLIDES = 27

HERE = Path(__file__).resolve().parent
IMG = HERE.parent / "docs" / "images"


# ──────────────────────────────────────────────────────────────
# 헬퍼
# ──────────────────────────────────────────────────────────────

def set_background(slide, color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_image(slide, path, left, top, width, height):
    if Path(path).exists():
        return slide.shapes.add_picture(
            str(path), Inches(left), Inches(top), Inches(width), Inches(height)
        )
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_LIGHT
    box.line.color.rgb = GRAY_MID
    return box


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.15), prs.slide_width, Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.22), Inches(12.2), Inches(0.85)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY_LIGHT


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=GRAY_DARK, align=PP_ALIGN.LEFT,
             line_spacing=1.25, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.text = line
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
        if not p.runs:
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.italic = italic
            p.font.color.rgb = color
    return tb


def add_bullets(slide, bullets, left, top, width, height,
                size=16, color=GRAY_DARK, line_spacing=1.35, bullet="•"):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  " + b
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return tb


def add_footer(slide, page_no):
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.35)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = (
        f"Police Frameworks · /peel · 최희철 "
        f"(경주경찰서 경찰발전협의회 회원, 소프트웨어 개발자) · {page_no}/{TOTAL_SLIDES}"
    )
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY_MID


def add_highlight_box(slide, text, left, top, width, height, fill=YELLOW_HL,
                      size=16, border=GOLD):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.2)
    box.adjustments[0] = 0.12
    tb = slide.shapes.add_textbox(
        Inches(left + 0.25), Inches(top + 0.16),
        Inches(width - 0.5), Inches(height - 0.32),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = GRAY_DARK
        p.line_spacing = 1.3


# ──────────────────────────────────────────────────────────────
# 사례 데이터 (5건, 각 3장)
# ──────────────────────────────────────────────────────────────

CASES = [
    # ═══════════════════════════════════════════════════════════
    # 사례 1 · 성건동 공원 청소년
    # ═══════════════════════════════════════════════════════════
    {
        "num": 1,
        "title": "성건동 공원",
        "subtitle": "3개월째 같은 민원, 어떻게 풀까",
        "image": "case1_park.png",
        "image_side": "left",
        "situation_title": "이런 상황 있으시죠?",
        "situation": (
            "경주시 성건동 ○○근린공원, 밤 10시부터 새벽 2시.\n"
            "고등학생 대여섯 명이 술·담배·소란.\n"
            "순찰차가 가면 흩어지고, 떠나면 다시 복귀.\n"
            "민원은 3개월째 14건 쌓였습니다.\n\n"
            "한쪽 주민: \"강하게 단속해라\"\n"
            "한쪽 주민: \"우리 애들 낙인찍지 마라\"\n"
            "성건파출소 야간 인력은 상시 배치 불가."
        ),
        "hook": (
            "단속만으로는 옆 동네 공원으로 옮겨갈 뿐.\n"
            "다음 장에서 /peel 이 이 민원을 어떻게 쪼개는지 보여드립니다."
        ),
        "selected": [
            ("SARA", "4단계 진단 사이클"),
            ("Crime Triangle", "6박스 원인 분해"),
            ("Hot Spots", "시간·공간 집중"),
            ("CPTED", "조명·조경·시야"),
            ("COP", "학교·센터·공원녹지과"),
            ("Procedural Justice", "청소년·주민 존엄"),
            ("Broken Windows 비판", "\"무관용\" 함정 필터"),
        ],
        "excluded": [
            ("PEACE·Cognitive", "조사 단계 아님"),
            ("BCSM·ICAT", "현장 위기 아님"),
            ("ILP", "단일 지점 전술 문제"),
        ],
        "pipeline": [
            "Scanning\n14건 지도화",
            "Analysis\n6박스 분해",
            "BW 필터\n단속 함정",
            "Response\nCPTED+COP",
            "PJ 삽입\n청소년 존엄",
            "Assessment\n3개월 재측정",
        ],
        "c_subtitle": "7개 공구가 실제로 어떻게 움직이나",
        "c_actions": [
            {
                "title": "진단 TF — 어디가 진짜 문제인가",
                "frameworks": ["SARA", "Crime Triangle", "Hot Spots"],
                "owner": "성건파출소 + 경주서 범죄예방진단팀",
                "deadline": "1주",
                "steps": [
                    "KICS로 최근 3개월 ○○근린공원 반경 112 신고 전수 추출",
                    "신고 14건을 지도에 점찍어 북측 놀이터·정자 집중 확인",
                    "Crime Triangle 6박스 워크시트로 Handler·Place·Manager 개입점 특정",
                ],
            },
            {
                "title": "CPTED 환경 개입 — 조명과 시야부터",
                "frameworks": ["CPTED"],
                "owner": "경주서 생활안전과 + 경주시 공원녹지과",
                "deadline": "1개월",
                "steps": [
                    "공원 야간 합동 점검 (조명 밝기·사각지대·정자 음영)",
                    "조명 LED 교체 예산 요청 공문 발송",
                    "조경 가지치기·CCTV 각도 조정",
                ],
            },
            {
                "title": "COP 파트너십 — 학교·센터·편의점",
                "frameworks": ["COP"],
                "owner": "성건파출소장 + 여성청소년과",
                "deadline": "2주",
                "steps": [
                    "인근 중·고교 생활지도부 연락망 구축",
                    "경주시 청소년상담복지센터 대안 활동 연계",
                    "주변 편의점주 3곳 방문, 청소년 주류 판매 점검",
                ],
            },
            {
                "title": "주민·청소년 대응 — 4요소 언어",
                "frameworks": ["Procedural Justice"],
                "owner": "성건파출소장",
                "deadline": "2주",
                "steps": [
                    "Voice·Neutrality·Respect·Trust 스크립트 마련",
                    "주민 설명회 개최 (단속 요구 주민 포함)",
                    "청소년 발견 시 개별 대화 프로토콜 전 경찰관 공유",
                ],
            },
            {
                "title": "\"무관용 단속\" 요구 필터",
                "frameworks": ["Broken Windows 비판"],
                "owner": "지구대장",
                "deadline": "즉시",
                "steps": [
                    "Zero Tolerance의 함정·낙인 효과 논리 정리",
                    "주민 설명회에서 \"왜 단속만으로 안 되는지\" 발표",
                    "경찰관 전원에게 낙인 방지 원칙 전달",
                ],
            },
            {
                "title": "3개월 후 재측정",
                "frameworks": ["SARA Assessment"],
                "owner": "경주서 생활안전계",
                "deadline": "3개월 후",
                "steps": [
                    "○○공원 야간 민원 건수 재측정",
                    "인근 지역 범죄 이동 확인",
                    "주민 간이 설문으로 체감 안전 측정",
                ],
            },
        ],
    },
    # ═══════════════════════════════════════════════════════════
    # 사례 2 · 황리단길 야간
    # ═══════════════════════════════════════════════════════════
    {
        "num": 2,
        "title": "황리단길 야간",
        "subtitle": "관광지가 된 뒤, 순찰 인력은 그대로",
        "image": "case2_hwangnidan.png",
        "image_side": "right",
        "situation_title": "이번엔 이런 상황",
        "situation": (
            "황남동 황리단길이 관광명소가 된 지 수년.\n"
            "주말 밤마다 관광객·상인·주민이 섞입니다.\n"
            "야간 주취·소란·주차 마찰·소음 신고.\n\n"
            "황남파출소 관할 인력은 그대로인데,\n"
            "유동 인구는 몇 배로 늘었습니다.\n"
            "전부 순찰로 막을 수 있는 규모가 아닙니다.\n"
            "상인은 영업을, 주민은 수면을 요구합니다."
        ),
        "hook": (
            "공간은 좁고 인력은 한정.\n"
            "다음 장에서 /peel 이 어떻게 \"시간·지점·파트너\"로 쪼개는지 보여드립니다."
        ),
        "selected": [
            ("Hot Spots", "22~02시 × 특정 블록"),
            ("ILP", "데이터로 우선순위"),
            ("COP", "상인협·주민자치위"),
            ("Procedural Justice", "4요소 접촉 훈련"),
            ("Broken Windows 비판", "\"엄격 단속\" 검토"),
        ],
        "excluded": [
            ("PEACE·Cognitive", "조사 면담 단계 아님"),
            ("BCSM·ICAT", "현장 위기 아님"),
            ("CPTED", "관광지 설계 변경 아님"),
            ("SARA 독자", "Hot Spots 우선"),
        ],
        "pipeline": [
            "112 분석\n시간·지점",
            "Koper 15분\n순찰 전략",
            "상인·주민\n협의체",
            "PJ 4요소\n접촉 훈련",
            "분기 평가\n체감·신고",
        ],
        "c_subtitle": "좁은 공간·한정 인력을 어떻게 쓰는가",
        "c_actions": [
            {
                "title": "Hot Spots + ILP 데이터 분석",
                "frameworks": ["Hot Spots", "ILP"],
                "owner": "황남파출소 + 경주서 범죄예방진단팀",
                "deadline": "2주",
                "steps": [
                    "최근 6개월 주말 야간(22~02시) 112 신고 유형별 분포",
                    "핫스팟 블록 상위 3곳 선정",
                    "Koper 15분 규칙 기반 순찰 배치 계획 수립",
                ],
            },
            {
                "title": "COP 상인·주민 협의체",
                "frameworks": ["COP"],
                "owner": "황남파출소장",
                "deadline": "1개월",
                "steps": [
                    "상인협의회 + 주민자치위 첫 합동 간담회 주재",
                    "공동 이슈 목록 수렴 (주취·주차·소음)",
                    "월 1회 정기 회의 체계 합의",
                ],
            },
            {
                "title": "Procedural Justice 접촉 훈련",
                "frameworks": ["Procedural Justice"],
                "owner": "황남파출소 + 교육 담당",
                "deadline": "3주",
                "steps": [
                    "관광객·상인·주민별 접촉 스크립트 마련 (4요소 삽입)",
                    "지구대 전원 롤플레이 훈련",
                    "피드백 수렴 체계 구축",
                ],
            },
            {
                "title": "\"엄격 단속\" 요구 필터",
                "frameworks": ["Broken Windows 비판"],
                "owner": "지구대장",
                "deadline": "즉시",
                "steps": [
                    "관광지 Zero Tolerance 실패 사례 공유",
                    "상인 대립·주민 대립 완화 메시지 정리",
                    "대안(협의체·핫스팟·PJ)을 \"단속 강화\"보다 앞에 제시",
                ],
            },
            {
                "title": "관광 성수기 특별 운영",
                "frameworks": ["Hot Spots", "COP"],
                "owner": "경주서 생활안전계",
                "deadline": "성수기 2주 전",
                "steps": [
                    "봄·가을 성수기 주말 전용 인력 편성",
                    "자원봉사·지역경찰 혼합 운영",
                    "경주시 관광과와 공동 대응 매뉴얼",
                ],
            },
            {
                "title": "분기별 평가",
                "frameworks": ["SARA Assessment"],
                "owner": "경주서 생활안전계",
                "deadline": "분기",
                "steps": [
                    "민원 건수·상인 만족도·주민 체감 측정",
                    "인근 지역 범죄 이동 확인",
                    "차기 분기 우선순위 조정",
                ],
            },
        ],
    },
    # ═══════════════════════════════════════════════════════════
    # 사례 3 · 외동산단 새벽 절도
    # ═══════════════════════════════════════════════════════════
    {
        "num": 3,
        "title": "외동산단 새벽 절도",
        "subtitle": "순찰보다 조명이 더 싸고 강하다",
        "image": "case3_industrial.png",
        "image_side": "left",
        "situation_title": "또 다른 흔한 상황",
        "situation": (
            "외동읍 산업단지 공장 주차장.\n"
            "새벽 시간 차량 내 물품 절도 12건 / 3개월.\n"
            "지금까지 대응: 순찰 강화.\n"
            "결과: 순찰 빠지면 다시 발생.\n\n"
            "외동파출소 새벽 인력은 한계.\n"
            "단지 관리공단은 방범이 주 업무가 아닙니다.\n"
            "입주기업은 각자 보안을 각자 책임집니다."
        ),
        "hook": (
            "Sherman·Weisburd 연구: 핫스팟 15분 체류만으로 효과 급증.\n"
            "다음 장: /peel 이 순찰 강화 대신 CPTED를 왜 먼저 권하는지."
        ),
        "selected": [
            ("SARA", "문제 구조화"),
            ("Crime Triangle", "가해자·대상·장소"),
            ("Hot Spots", "구역 집중 확인"),
            ("CPTED", "조명·CCTV·차단기"),
            ("COP", "관리공단·입주기업"),
        ],
        "excluded": [
            ("PEACE·Cognitive", "조사 단계 아님"),
            ("BCSM·ICAT", "현장 위기 아님"),
            ("Procedural Justice", "주 개입 아님 (전제)"),
            ("Broken Windows", "무질서 프레임 부적절"),
        ],
        "pipeline": [
            "신고 정리\n12건 지도화",
            "6박스\n분해",
            "CPTED\n현장 점검",
            "Response\n조명·CCTV",
            "관리공단\n책임 분담",
            "3개월\n재측정",
        ],
        "c_subtitle": "순찰이 아닌 환경·파트너로 푸는 법",
        "c_actions": [
            {
                "title": "SARA + Crime Triangle 진단",
                "frameworks": ["SARA", "Crime Triangle"],
                "owner": "외동파출소 + 경주서 범죄예방진단팀",
                "deadline": "1주",
                "steps": [
                    "12건 신고 시공간 분포·수법·피해 품목 집계",
                    "6박스 분해로 Place·Manager 개입점 도출",
                    "공통 피해 패턴(차종·보관물·시간대) 문서화",
                ],
            },
            {
                "title": "Hot Spots 순찰 재설계",
                "frameworks": ["Hot Spots"],
                "owner": "외동파출소",
                "deadline": "2주",
                "steps": [
                    "상위 5% 핫스팟 구역 특정",
                    "새벽 시간대 Koper 15분 집중 순찰",
                    "순찰 로그 정기 공유 (주간 단위)",
                ],
            },
            {
                "title": "CPTED 합동 현장 점검",
                "frameworks": ["CPTED"],
                "owner": "경주서 + 단지 관리공단",
                "deadline": "1개월",
                "steps": [
                    "주차장 조명 수명·밝기 점검 (20m 얼굴 식별 기준)",
                    "CCTV 각도·사각지대 매핑",
                    "차단기·출입 통제 작동 확인, 개선 우선순위 작성",
                ],
            },
            {
                "title": "COP 관리공단·기업 협력",
                "frameworks": ["COP"],
                "owner": "외동파출소장",
                "deadline": "2주",
                "steps": [
                    "단지 관리공단과 방범 MOU 추진",
                    "입주기업 방범 담당자 연락망 구축",
                    "월 1회 브리핑 체계 (사건 공유·예방 팁)",
                ],
            },
            {
                "title": "환경 개선 예산 확보",
                "frameworks": ["CPTED", "COP"],
                "owner": "경주서 생활안전과",
                "deadline": "2개월",
                "steps": [
                    "조명·CCTV 개선 견적 3곳 확보",
                    "관리공단·경주시·경북경찰청 예산 협의",
                    "시공 일정 확정 및 공유",
                ],
            },
            {
                "title": "3개월 재측정",
                "frameworks": ["SARA Assessment"],
                "owner": "경주서 생활안전계",
                "deadline": "3개월 후",
                "steps": [
                    "신고 건수·피해액 재측정",
                    "인접 단지 범죄 이동 확인",
                    "SARA 사이클 재시작 여부 결정",
                ],
            },
        ],
    },
    # ═══════════════════════════════════════════════════════════
    # 사례 4 · 고압 민원
    # ═══════════════════════════════════════════════════════════
    {
        "num": 4,
        "title": "규정대로 했는데 민원",
        "subtitle": "30년 경력이 가장 억울한 순간",
        "image": "case4_station.png",
        "image_side": "right",
        "situation_title": "이게 제일 억울하지 않으십니까",
        "situation": (
            "지구대에서 주취자를 규정대로 처리.\n"
            "며칠 뒤 \"경찰이 고압적이었다\" 민원.\n\n"
            "내부적으로는 절차 다 지켰습니다.\n"
            "그런데 왜 민원이 나왔을까요?\n\n"
            "Tom Tyler(미국 사회심리학자) 연구:\n"
            "시민은 \"결과\"보다 \"과정\"을 기억합니다.\n"
            "4요소 중 하나라도 부족하면 신뢰가 깨집니다."
        ),
        "hook": (
            "\"규정대로\"는 정당성 언어가 아닙니다.\n"
            "다음 장에서 /peel 이 4요소로 사건을 해부하는 법을 보여드립니다."
        ),
        "selected": [
            ("Procedural Justice", "4요소 사건 해부"),
            ("COP", "지구대장 직접 연락"),
            ("SARA", "재발 방지 구조 분석"),
        ],
        "excluded": [
            ("Hot Spots·CPTED", "장소 문제 아님"),
            ("PEACE·Cognitive", "이미 종결"),
            ("BCSM·ICAT", "위기 아님"),
            ("Crime Triangle", "범죄 분석 아님"),
            ("ILP", "데이터 규모 아님"),
            ("Broken Windows", "무관"),
        ],
        "pipeline": [
            "4요소\n사건 평가",
            "Voice\n주민 청취",
            "Respect\n사과 필요 시",
            "SARA\n재발 방지",
            "교육\n언어·태도",
        ],
        "c_subtitle": "규정을 지켰어도 신뢰가 깨진 이유를 푼다",
        "c_actions": [
            {
                "title": "4요소로 사건 해부",
                "frameworks": ["Procedural Justice"],
                "owner": "감찰 + 생활안전계",
                "deadline": "1주",
                "steps": [
                    "Voice·Neutrality·Respect·Trust 항목별 당일 대응 점검",
                    "녹화·녹음 재검토, 4요소 중 부족한 지점 특정",
                    "1페이지 분석 보고서 작성",
                ],
            },
            {
                "title": "지구대장 주민 직접 연락",
                "frameworks": ["COP", "Procedural Justice"],
                "owner": "지구대장",
                "deadline": "3일 내",
                "steps": [
                    "민원인에게 직접 전화 또는 방문",
                    "4요소 언어로 사건을 다시 설명",
                    "필요 시 사과, 후속 조치 약속",
                ],
            },
            {
                "title": "유사 민원 빈도 확인",
                "frameworks": ["SARA", "ILP"],
                "owner": "감찰",
                "deadline": "2주",
                "steps": [
                    "지난 1년 유사 유형 민원 건수 집계",
                    "특정 시간대·경찰관 편중 여부 확인",
                    "패턴이면 구조적 원인 진단",
                ],
            },
            {
                "title": "접촉 표준 스크립트 개선",
                "frameworks": ["Procedural Justice"],
                "owner": "경주서 교육과",
                "deadline": "3주",
                "steps": [
                    "주취자 처리 표준 스크립트 재작성 (4요소 삽입)",
                    "지구대 롤플레이 훈련",
                    "분기 평가 포함",
                ],
            },
            {
                "title": "재발 방지 교육",
                "frameworks": ["Procedural Justice", "COP"],
                "owner": "경주서 교육과",
                "deadline": "1개월",
                "steps": [
                    "Procedural Justice 기본 교육 실시",
                    "본 사건 포함 사례 중심 워크숍",
                    "교육 효과 평가 및 기록",
                ],
            },
        ],
    },
    # ═══════════════════════════════════════════════════════════
    # 사례 5 · 거래처 사칭 사기
    # ═══════════════════════════════════════════════════════════
    {
        "num": 5,
        "title": "거래처 사칭 사기",
        "subtitle": "순찰로는 절대 못 잡는 범죄",
        "image": "case5_fraud.png",
        "image_side": "left",
        "situation_title": "경제범죄도 /peel 이 다룹니다",
        "situation": (
            "거래처 구매 담당자를 사칭해\n"
            "\"공사 건이 급해서 자재를 먼저 사서 보내달라\"\n"
            "피해업체가 결제 후 연락 두절.\n\n"
            "카카오톡 프로필 복제, 이메일 도메인 변조\n"
            "(samsung → samsnug), 발신번호 변작.\n\n"
            "단일 사건이 아닌 조직·연속·디지털 범죄.\n"
            "물리적 공간이 없어 순찰로 막을 수 없습니다."
        ),
        "hook": (
            "이 범죄는 순찰로 못 잡습니다.\n"
            "다음 장: /peel 이 데이터·예방·진술·면담으로 어떻게 엮는지."
        ),
        "selected": [
            ("SARA", "문제 정의 구조"),
            ("ILP", "계좌·전화·IP 분석"),
            ("Crime Triangle", "통로 차단 지점"),
            ("COP", "조달담당자 교육"),
            ("Cognitive Interview", "진술 단서 인출"),
            ("PEACE", "검거 시 면담"),
        ],
        "excluded": [
            ("CPTED·Hot Spots", "물리 공간 범죄 아님"),
            ("BCSM·ICAT", "위기 대응 아님"),
            ("Broken Windows", "해당 없음"),
            ("Procedural Justice", "기본 전제 (2차 가해 금지)"),
        ],
        "pipeline": [
            "SARA\n유형 집계",
            "ILP\n계좌·전화\n클러스터링",
            "Triangle\n통로 차단",
            "COP\n조달 교육",
            "Cognitive\n진술 인출",
            "PEACE\n검거 면담",
        ],
        "c_subtitle": "데이터 → 예방 → 진술 → 면담의 4박자",
        "c_actions": [
            {
                "title": "SARA + ILP 합동 TF 결성",
                "frameworks": ["SARA", "ILP"],
                "owner": "경주서 생활안전계 + 사이버수사대",
                "deadline": "1주",
                "steps": [
                    "최근 3개월 \"거래처 사칭\" 신고 집계 (건수·피해액·업종)",
                    "동일 계좌·전화번호 연계사건 탐지",
                    "텍스트 클러스터링 (문구·수법·말투 패턴)",
                ],
            },
            {
                "title": "Crime Triangle 통로 차단",
                "frameworks": ["Crime Triangle"],
                "owner": "경주서 + 본청 사이버수사국",
                "deadline": "3주",
                "steps": [
                    "카카오·통신사 플랫폼 협조 요청 (프로필 복제·발신 변작)",
                    "금융정보분석원에 대포계좌 동결 요청",
                    "국제공조 필요 시 검경 연계 조기 가동",
                ],
            },
            {
                "title": "COP 조달담당자 실전 교육",
                "frameworks": ["COP"],
                "owner": "경주서 생활안전·범죄예방과",
                "deadline": "3주",
                "steps": [
                    "상공회의소·산업단지 관리공단과 MOU 체결",
                    "실제 사칭 캡처 포함 교육 자료 개발",
                    "\"신규 계좌 입금 전 2인 재확인\" 프로토콜 배포",
                ],
            },
            {
                "title": "조기 경보 네트워크",
                "frameworks": ["COP"],
                "owner": "경주서 생활안전과",
                "deadline": "1개월",
                "steps": [
                    "상공회의소 경보 문자 발송 채널 구축",
                    "피해 발생 48시간 내 인근 업체 전파 체계",
                    "자율방범대·지역 밴드·카페 활용",
                ],
            },
            {
                "title": "Cognitive Interview 진술 체계",
                "frameworks": ["Cognitive Interview"],
                "owner": "경주서 수사과 + 교육",
                "deadline": "1개월",
                "steps": [
                    "피해자 조서 양식에 4기법 반영 (문맥·세부·역순·관점)",
                    "조사관 교육 및 롤플레이",
                    "사칭범 말투·단어·시간대를 ILP로 피드백",
                ],
            },
            {
                "title": "PEACE 면담 준비",
                "frameworks": ["PEACE"],
                "owner": "경주서 수사과",
                "deadline": "검거 시",
                "steps": [
                    "증거(계좌·통신·진술) 완벽 정리 후 면담",
                    "Account → Clarification → Challenge 구조 적용",
                    "조직 전모 확인이 목표 (개인 자백 유도 아님)",
                ],
            },
        ],
    },
]


# ──────────────────────────────────────────────────────────────
# 사례 A (상황)
# ──────────────────────────────────────────────────────────────

def make_case_slide_a(prs, case, page_no):
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_title_bar(s, f"사례 {case['num']} · {case['title']}", case["subtitle"])

    img_path = IMG / case["image"]
    if case["image_side"] == "left":
        add_image(s, img_path, 0.5, 1.5, 6.0, 3.5)
        text_x = 6.85
    else:
        add_image(s, img_path, 6.85, 1.5, 6.0, 3.5)
        text_x = 0.6

    add_text(s, case["situation_title"], text_x, 1.5, 6.2, 0.5,
             size=20, bold=True, color=NAVY)
    add_text(s, case["situation"], text_x, 2.1, 6.2, 2.9,
             size=14, color=GRAY_DARK, line_spacing=1.35)

    add_highlight_box(s, case["hook"], 0.6, 5.3, 12.1, 1.55, size=15)
    add_footer(s, page_no)
    return s


# ──────────────────────────────────────────────────────────────
# 사례 B (/peel 분석 — 선정/제외/파이프라인)
# ──────────────────────────────────────────────────────────────

def make_case_slide_b(prs, case, page_no):
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_title_bar(s,
                  f"사례 {case['num']} · /peel 의 분석",
                  "어떤 공구를 꺼내고, 어떤 공구는 뺄지")

    # 선정 (녹색)
    sel_bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(6.05), Inches(0.4)
    )
    sel_bar.fill.solid()
    sel_bar.fill.fore_color.rgb = GREEN
    sel_bar.line.fill.background()
    add_text(s, "✓ 선정 — 이 공구들을 꺼냅니다",
             0.75, 1.55, 5.9, 0.32, size=13, bold=True, color=WHITE)

    y = 2.0
    for name, reason in case["selected"]:
        add_text(s, f"•  {name}", 0.75, y, 3.2, 0.3,
                 size=13, bold=True, color=GREEN)
        add_text(s, reason, 3.95, y + 0.02, 2.7, 0.3,
                 size=11, color=GRAY_MID)
        y += 0.32

    # 제외 (회색)
    ex_bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(1.5), Inches(6.05), Inches(0.4)
    )
    ex_bar.fill.solid()
    ex_bar.fill.fore_color.rgb = GRAY_MID
    ex_bar.line.fill.background()
    add_text(s, "✕ 제외 — 이 공구들은 뺐습니다",
             7.0, 1.55, 5.9, 0.32, size=13, bold=True, color=WHITE)

    y = 2.0
    for name, reason in case["excluded"]:
        add_text(s, f"•  {name}", 7.0, y, 3.2, 0.3,
                 size=13, bold=True, color=GRAY_MID)
        add_text(s, reason, 10.2, y + 0.02, 2.7, 0.3,
                 size=11, color=GRAY_MID)
        y += 0.32

    # 파이프라인
    add_text(s, "→ 이 순서로 진행합니다",
             0.6, 4.55, 12.1, 0.3, size=14, bold=True, color=NAVY)

    steps = case["pipeline"]
    n = len(steps)
    total_w = 12.1
    gap = 0.15
    box_w = (total_w - gap * (n - 1)) / n
    x = 0.6
    for i, step in enumerate(steps):
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(4.95), Inches(box_w), Inches(1.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.adjustments[0] = 0.15
        add_text(s, step, x + 0.05, 5.08, box_w - 0.1, 0.95,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)
        if i < n - 1:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + box_w + 0.015), Inches(5.37), Inches(0.12), Inches(0.25)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = GOLD
            arr.line.fill.background()
        x += box_w + gap

    add_highlight_box(
        s,
        "다음 장에서 각 공구가 현장에서 구체적으로 무엇을 하는지 정리해 드립니다.",
        0.6, 6.3, 12.1, 0.75, size=13
    )

    add_footer(s, page_no)
    return s


# ──────────────────────────────────────────────────────────────
# 사례 C (실행 방안 상세 — 6개 카드)
# ──────────────────────────────────────────────────────────────

def make_case_slide_c(prs, case, page_no):
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_title_bar(s,
                  f"사례 {case['num']} · 구체적인 실행 방안",
                  case.get("c_subtitle", ""))

    actions = case["c_actions"][:6]

    # 2열 × 3행
    positions = [
        (0.4, 1.35), (6.85, 1.35),
        (0.4, 3.33), (6.85, 3.33),
        (0.4, 5.31), (6.85, 5.31),
    ]
    card_w = 6.05
    card_h = 1.85

    for i, action in enumerate(actions):
        if i >= len(positions):
            break
        x, y = positions[i]

        # 카드 박스
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(card_w), Inches(card_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = NAVY
        box.line.width = Pt(1.3)
        box.adjustments[0] = 0.1

        # 좌측 번호 스트립
        strip = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(0.5), Inches(card_h)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = NAVY
        strip.line.fill.background()
        add_text(s, str(i + 1), x, y + 0.55, 0.5, 0.6,
                 size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

        # 제목
        add_text(s, action["title"], x + 0.6, y + 0.1, card_w - 0.75, 0.35,
                 size=13, bold=True, color=NAVY)

        # 프레임워크 태그
        fw_text = "  ·  ".join(action["frameworks"])
        add_text(s, fw_text, x + 0.6, y + 0.43, card_w - 0.75, 0.25,
                 size=9, bold=True, color=GOLD, italic=True)

        # 담당 · 기한
        meta = f"{action['owner']}   |   기한 {action['deadline']}"
        add_text(s, meta, x + 0.6, y + 0.67, card_w - 0.75, 0.25,
                 size=9, color=GRAY_MID)

        # 단계 (3개)
        tb = s.shapes.add_textbox(
            Inches(x + 0.6), Inches(y + 0.92),
            Inches(card_w - 0.75), Inches(card_h - 1.0)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        for j, step in enumerate(action["steps"][:3]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "•  " + step
            p.line_spacing = 1.2
            p.space_after = Pt(0)
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(9)
                run.font.color.rgb = GRAY_DARK

    add_footer(s, page_no)
    return s


# ══════════════════════════════════════════════════════════════
# 프레젠테이션 생성
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════
# Slide 1 — 표지
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, NAVY_DEEP)
add_image(s, IMG / "cover.png", 0, 0, 13.333, 7.5)

overlay = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(3.7), prs.slide_width, Inches(3.8)
)
overlay.fill.solid()
overlay.fill.fore_color.rgb = NAVY_DEEP
overlay.line.fill.background()

add_text(s, "POLICE FRAMEWORKS", 0.8, 0.6, 12, 0.4,
         size=14, bold=True, color=GOLD)
add_text(s, "3개월째 똑같은 민원, 이번엔 순서부터",
         0.8, 4.0, 12, 1.0, size=40, bold=True, color=WHITE)
add_text(s, "경찰 프레임워크 툴킷  /  메타 라우터 /peel",
         0.8, 4.95, 12, 0.55, size=22, color=GOLD)
add_text(s,
         "12개 근거기반 경찰활동 프레임워크를, 상황에 맞게 AI가 골라드립니다",
         0.8, 5.6, 12, 0.5, size=16, color=GRAY_LIGHT)

line = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.3), Inches(2.0), Inches(0.06)
)
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

add_text(s, "최희철", 0.8, 6.45, 12, 0.4, size=18, bold=True, color=WHITE)
add_text(s,
         "경주경찰서 경찰발전협의회 회원  ·  소프트웨어 개발자  ·  2026년 4월",
         0.8, 6.78, 12, 0.35, size=13, color=GRAY_LIGHT)


# ══════════════════════════════════════════════════════════════
# Slide 2 — 차례
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "오늘 이야기 순서", "사례 먼저 · 이론은 짧게 · 설치는 천천히")

add_text(s, "결론부터 말씀드립니다.",
         0.6, 1.6, 12, 0.5, size=20, bold=True, color=NAVY)
add_text(s,
         "30년 현장 경험을 부정하는 도구가 아닙니다.\n"
         "여러분이 이미 머리로 하시는 일을 한 장의 수첩처럼 정리해 드리는 도구입니다.",
         0.6, 2.15, 12, 1.2, size=15, color=GRAY_MID)

stages = [
    ("1", "사례 5가지", "경주 현장에서 바로\n있을 만한 이야기\n(사례당 3장씩)", RED),
    ("2", "이론은 간단히", "12개 '공구'가 뭔지\n한눈에 보여드립니다\n(6장)", NAVY),
    ("3", "설치는 천천히", "Claude에게 링크 주고\n\"깔아줘\" 한마디\n(3장)", GOLD),
]
x = 0.6
for num, title, desc, color in stages:
    circle = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(3.8), Inches(1.0), Inches(1.0)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(s, num, x + 0.3, 3.92, 1.0, 0.8,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(5.1), Inches(4.0), Inches(1.85)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.adjustments[0] = 0.1
    add_text(s, title, x + 0.2, 5.22, 3.6, 0.5,
             size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.2, 5.8, 3.6, 1.1,
             size=13, color=GRAY_MID, align=PP_ALIGN.CENTER)
    x += 4.3

add_footer(s, 2)


# ══════════════════════════════════════════════════════════════
# Slides 3-17 — 사례 5건 × 3장
# ══════════════════════════════════════════════════════════════
page = 3
for case in CASES:
    make_case_slide_a(prs, case, page)
    make_case_slide_b(prs, case, page + 1)
    make_case_slide_c(prs, case, page + 2)
    page += 3


# ══════════════════════════════════════════════════════════════
# Slide 18 — 공구함 비유
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "그래서 이게 뭔가요?", "30년 수첩을 한 장으로, 공구함과 공구처럼")

add_image(s, IMG / "toolbox.png", 6.85, 1.5, 6.0, 4.5)

add_text(s, "비유 — 고참 반장님의 수첩", 0.6, 1.6, 6.2, 0.6,
         size=20, bold=True, color=NAVY)
body = (
    "30년 현장을 뛴 반장이 \"이런 민원엔 이렇게\"\n"
    "라고 적어둔 노트가 있다고 합시다.\n\n"
    "이 도구는 그 노트를 — 세계 여러 나라 경찰 연구로\n"
    "검증된 — 12개의 공구로 정리한 것입니다.\n\n"
    "그리고 /peel 이라는 안내 데스크가\n"
    "상황에 맞는 공구를 골라드립니다.\n\n"
    "볼트엔 렌치, 나사엔 드라이버.\n"
    "민원엔 절차적 정의, 반복엔 SARA."
)
add_text(s, body, 0.6, 2.3, 6.2, 4.7,
         size=15, color=GRAY_DARK, line_spacing=1.4)

add_footer(s, 18)


# ══════════════════════════════════════════════════════════════
# Slide 19 — 12개 공구 4개 서랍
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "12개 공구, 4개 서랍", "상황에 맞게 골라 쓰는 공구함")

drawers = [
    ("문제 패턴을 분석할 때",
     ["SARA", "범죄삼각형", "CPTED", "핫스팟", "ILP"],
     "반복 민원·핫스팟·환경 개선",
     NAVY),
    ("조사 · 면담할 때",
     ["PEACE 모델", "인지 면담"],
     "피의자·피해자 진술, 강압 없는 면담",
     NAVY_LIGHT),
    ("현장 위기 대응할 때",
     ["BCSM", "ICAT"],
     "자살·인질·정신건강 위기, 물리력 최소화",
     RED),
    ("신뢰 · 지역사회",
     ["절차적 정의", "COP", "깨진유리창(비판)"],
     "민원 응대, 주민 관계, \"단속 요구\" 검토",
     GOLD),
]

CHIP_AREA_X = 4.7
CHIP_AREA_END = 12.95
CHIP_AREA_W = CHIP_AREA_END - CHIP_AREA_X
CHIP_GAP = 0.12

y = 1.55
for title, tools, desc, color in drawers:
    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y), Inches(4.0), Inches(1.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(s, title, 0.75, y + 0.22, 3.8, 0.45,
             size=16, bold=True, color=WHITE)
    add_text(s, desc, 0.75, y + 0.74, 3.8, 0.5,
             size=11, color=GRAY_LIGHT)

    n = len(tools)
    chip_w = (CHIP_AREA_W - CHIP_GAP * (n - 1)) / n
    tx = CHIP_AREA_X
    chip_font = 11 if n >= 5 else (14 if n == 3 else 17)

    for t in tools:
        chip = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(tx), Inches(y + 0.3), Inches(chip_w), Inches(0.7),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = WHITE
        chip.line.color.rgb = color
        chip.line.width = Pt(1.8)
        chip.adjustments[0] = 0.3
        add_text(s, t, tx + 0.05, y + 0.45, chip_w - 0.1, 0.4,
                 size=chip_font, bold=True, color=color, align=PP_ALIGN.CENTER)
        tx += chip_w + CHIP_GAP

    y += 1.42

add_footer(s, 19)


# ══════════════════════════════════════════════════════════════
# Slide 20 — 서랍 1 상세
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "서랍 1 · 문제 패턴 분석", "반복되는 민원의 뿌리를 찾을 때")

tools = [
    ("SARA",
     "탐색 → 분석 → 대응 → 평가 4단계.\n반복 민원을 근본부터 푸는 기본 순서표."),
    ("범죄 삼각형",
     "가해자·대상·장소 세 꼭짓점 + 각 억제자. 어디를 움직이면 가장 효과적인지 6박스로 본다."),
    ("CPTED",
     "환경 설계로 범죄 기회 제거. 조명·시야·관리 주체·활동 유치 6원칙."),
    ("핫스팟 경찰활동",
     "범죄는 관할 전역이 아니라 몇 블록에 집중. Koper 15분 규칙."),
    ("ILP",
     "데이터로 우선순위 정하기. 상습 가해자·반복 피해자·핫스팟 3대 대상."),
]

y = 1.55
for name, desc in tools:
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(y), Inches(12.1), Inches(1.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.2)
    box.adjustments[0] = 0.15
    name_bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(y), Inches(2.8), Inches(1.0),
    )
    name_bar.fill.solid()
    name_bar.fill.fore_color.rgb = NAVY
    name_bar.line.fill.background()
    add_text(s, name, 0.78, y + 0.3, 2.6, 0.5,
             size=16, bold=True, color=WHITE)
    add_text(s, desc, 3.6, y + 0.18, 9.0, 0.8,
             size=13, color=GRAY_DARK, line_spacing=1.3)
    y += 1.08

add_footer(s, 20)


# ══════════════════════════════════════════════════════════════
# Slide 21 — 서랍 2·3 상세
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "서랍 2·3 · 조사와 위기", "사람과 마주 앉을 때, 현장이 급할 때")

add_text(s, "서랍 2 — 조사·면담", 0.6, 1.55, 12, 0.4,
         size=17, bold=True, color=NAVY_LIGHT)
box = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.9),
)
box.fill.solid()
box.fill.fore_color.rgb = WHITE
box.line.color.rgb = NAVY_LIGHT
box.line.width = Pt(1.5)
box.adjustments[0] = 0.06

add_text(s, "PEACE 모델  ·  1992년 영국, 강압 없는 면담 5단계",
         0.85, 2.1, 11.8, 0.4, size=15, bold=True, color=NAVY)
add_text(s,
         "Planning → Engage → Account → Closure → Evaluate. 자백 압박이 아닌 진실 수집.\n"
         "허위자백 스캔들(Guildford Four 등) 이후 Reid 기법을 대체한 윤리적 모델.",
         0.85, 2.48, 11.8, 0.7, size=12, color=GRAY_MID, line_spacing=1.3)
add_text(s, "인지 면담  ·  피해자·목격자 기억 인출 극대화",
         0.85, 3.2, 11.8, 0.4, size=15, bold=True, color=NAVY)
add_text(s,
         "문맥 복원·모든 세부·역순·관점 전환 4기법. 메타분석 25~85% 더 많은 정확한 정보.",
         0.85, 3.55, 11.8, 0.3, size=12, color=GRAY_MID)

add_text(s, "서랍 3 — 현장 위기", 0.6, 4.15, 12, 0.4,
         size=17, bold=True, color=RED)
box2 = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.1),
)
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = RED
box2.line.width = Pt(1.5)
box2.adjustments[0] = 0.06

add_text(s, "BCSM  ·  FBI 위기협상 5단계 계단",
         0.85, 4.7, 11.8, 0.4, size=15, bold=True, color=RED)
add_text(s,
         "경청 → 공감 → 신뢰 → 영향 → 행동 변화. 계단은 건너뛸 수 없음.\n자살·인질·농성 협상.",
         0.85, 5.05, 11.8, 0.7, size=12, color=GRAY_MID, line_spacing=1.3)
add_text(s, "ICAT  ·  PERF 디에스컬레이션 훈련",
         0.85, 5.8, 11.8, 0.4, size=15, bold=True, color=RED)
add_text(s,
         "시간·거리·엄폐. 루이빌 실험: 물리력 28% ↓, 시민 부상 26% ↓, 경찰관 부상 36% ↓.",
         0.85, 6.15, 11.8, 0.4, size=12, color=GRAY_MID)

add_footer(s, 21)


# ══════════════════════════════════════════════════════════════
# Slide 22 — 서랍 4 상세
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "서랍 4 · 정당성과 지역사회", "규정만으로는 부족할 때")

add_text(s, "절차적 정의 — 네 가지를 다 지키셔야 신뢰가 쌓입니다",
         0.6, 1.55, 12, 0.4, size=16, bold=True, color=NAVY)

pj = [
    ("Voice", "발언권", "시민이 자기 입장을\n충분히 말할 기회"),
    ("Neutrality", "중립성", "사실과 규칙에 기반한\n일관된 판단"),
    ("Respect", "존중", "존엄을 침해하지 않는\n언어와 태도"),
    ("Trust", "신뢰", "선의로 시민 편에 선다는\n인상을 주기"),
]
x = 0.6
for en, kr, desc in pj:
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(2.0), Inches(3.0), Inches(2.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY
    card.line.color.rgb = GOLD
    card.line.width = Pt(2)
    card.adjustments[0] = 0.1
    add_text(s, en, x + 0.1, 2.15, 2.8, 0.4,
             size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, kr, x + 0.1, 2.5, 2.8, 0.5,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.1, 3.1, 2.8, 0.9,
             size=12, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)
    x += 3.15

box = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(4.3), Inches(6.0), Inches(2.3)
)
box.fill.solid()
box.fill.fore_color.rgb = WHITE
box.line.color.rgb = GOLD
box.line.width = Pt(1.5)
box.adjustments[0] = 0.06
add_text(s, "COP  ·  지역사회 경찰활동",
         0.85, 4.45, 5.5, 0.4, size=16, bold=True, color=NAVY)
add_text(s,
         "1829년 Robert Peel의 9대 원칙으로 회귀.\n"
         "\"경찰은 시민이고 시민은 경찰이다.\"\n\n"
         "지구대 밀착, 주민 협력, 문제해결 3축.",
         0.85, 4.85, 5.5, 1.7, size=13, color=GRAY_DARK, line_spacing=1.3)

box2 = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.85), Inches(4.3), Inches(6.0), Inches(2.3)
)
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = RED
box2.line.width = Pt(1.5)
box2.adjustments[0] = 0.06
add_text(s, "깨진 유리창  ·  비판적으로 다시 읽기",
         7.1, 4.45, 5.5, 0.4, size=16, bold=True, color=RED)
add_text(s,
         "원래 이론이 \"무관용 단속\"으로 오해되면서\n"
         "뉴욕 Stop-and-Frisk 같은 실패 사례가 나옴.\n\n"
         "이 도구는 \"엄격 단속\" 요구가 들어올 때\n"
         "경고 필터 역할을 합니다.",
         7.1, 4.85, 5.5, 1.7, size=13, color=GRAY_DARK, line_spacing=1.3)

add_footer(s, 22)


# ══════════════════════════════════════════════════════════════
# Slide 23 — /peel 안내데스크
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "/peel 이 어떻게 고르나", "종합병원 안내 데스크와 같습니다")

add_text(s, "비유 — 종합병원 안내 데스크",
         0.6, 1.6, 12, 0.5, size=20, bold=True, color=NAVY)
add_text(s,
         "환자가 \"허리가 아파요\" 하면 안내 데스크가 정형외과·신경외과·재활의학과\n"
         "어디로 가야 할지 정해줍니다. /peel 도 똑같습니다.",
         0.6, 2.1, 12, 1.1, size=15, color=GRAY_MID, line_spacing=1.4)

stages = [
    ("1단계", "상황 분류",
     "반복 패턴?\n단일 사건?\n현장 위기?"),
    ("2단계", "공구 고르기",
     "12개 중 2~6개 조합\n(나머지는 뺌)"),
    ("3단계", "순서 잡기",
     "어느 공구를\n먼저 · 나중에"),
    ("결과", "실행 방안",
     "담당·기한·단계까지\n구체적으로 제시"),
]
x = 0.6
for label, title, desc in stages:
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(3.4), Inches(2.9), Inches(2.9)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.8)
    box.adjustments[0] = 0.08
    lbl = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(3.4), Inches(2.9), Inches(0.6)
    )
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = NAVY
    lbl.line.fill.background()
    add_text(s, label, x + 0.1, 3.52, 2.7, 0.4,
             size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.1, 4.15, 2.7, 0.5,
             size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.1, 4.75, 2.7, 1.4,
             size=12, color=GRAY_MID, align=PP_ALIGN.CENTER, line_spacing=1.35)
    if x < 9:
        arr = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(x + 2.95), Inches(4.6), Inches(0.22), Inches(0.45)
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = GOLD
        arr.line.fill.background()
    x += 3.15

add_highlight_box(
    s,
    "경찰 실무자 눈높이에서 자동으로 초안이 나옵니다.\n"
    "여러분이 그 초안을 현장 경험으로 고쳐서 쓰시면 됩니다.",
    0.6, 6.5, 12.1, 0.9, size=13
)

add_footer(s, 23)


# ══════════════════════════════════════════════════════════════
# Slide 24 — Claude Skills 어디에 깔리나
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "어디에 깔리나요?", "Claude 3가지 제품 어디서든 똑같이 작동합니다")

add_text(s,
         "좋은 소식 — 한 번 깔면 세 곳 모두에서 쓸 수 있습니다.",
         0.6, 1.6, 12, 0.5, size=18, bold=True, color=NAVY)

products = [
    ("Claude Desktop 앱",
     "권장",
     "노트북·컴퓨터에\n깔아두고 쓰는 프로그램",
     ["대화창에서 바로 사용", "설정에서 업로드", "60대에게 가장 쉬움"],
     GOLD),
    ("Claude 웹 채팅 (claude.ai)",
     "지원",
     "브라우저로 접속해서 쓰는\n웹사이트 버전",
     ["로그인 후 바로 사용", "Customize → Skills", "Desktop과 동기화"],
     NAVY),
    ("Claude Code (개발자용)",
     "자동",
     "개발자가 쓰는\n명령어 기반 도구",
     ["슬래시 명령 자동 인식", "플러그인 명령으로 설치", "기술자 도움 시 편리"],
     NAVY_LIGHT),
]

x = 0.6
for name, badge, desc, features, color in products:
    # 카드
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(2.3), Inches(4.1), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    card.adjustments[0] = 0.08

    # 상단 색 바
    top = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(2.3), Inches(4.1), Inches(0.85)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = color
    top.line.fill.background()

    add_text(s, name, x + 0.15, 2.4, 3.8, 0.4,
             size=16, bold=True, color=WHITE)

    # 배지
    badge_shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.15), Inches(2.75), Inches(0.9), Inches(0.3)
    )
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = WHITE
    badge_shape.line.fill.background()
    badge_shape.adjustments[0] = 0.5
    add_text(s, badge, x + 0.15, 2.79, 0.9, 0.25,
             size=10, bold=True, color=color, align=PP_ALIGN.CENTER)

    # 설명
    add_text(s, desc, x + 0.2, 3.35, 3.7, 0.8,
             size=12, color=GRAY_DARK, line_spacing=1.3)

    # 특징 체크리스트
    y2 = 4.3
    for f in features:
        add_text(s, f"✓  {f}", x + 0.2, y2, 3.7, 0.4,
                 size=11, color=GRAY_MID)
        y2 += 0.4

    # 이 PC 표시 — Claude Code 카드에
    if name == "Claude Code (개발자용)":
        note = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.2), Inches(6.2), Inches(3.7), Inches(0.45)
        )
        note.fill.solid()
        note.fill.fore_color.rgb = YELLOW_HL
        note.line.color.rgb = GOLD
        note.line.width = Pt(1)
        note.adjustments[0] = 0.4
        add_text(s, "지금 이 노트북에 이미 설치됨", x + 0.25, 6.28, 3.6, 0.3,
                 size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)

    x += 4.3

add_footer(s, 24)


# ══════════════════════════════════════════════════════════════
# Slide 25 — Claude에게 설치 부탁하기
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "Claude에게 부탁하시면 됩니다",
              "링크 하나 주고 \"깔아줘\" 한마디")

add_text(s, "가장 쉬운 방법",
         0.6, 1.55, 12, 0.4, size=20, bold=True, color=NAVY)

# 큰 말풍선
bubble = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.8)
)
bubble.fill.solid()
bubble.fill.fore_color.rgb = YELLOW_HL
bubble.line.color.rgb = GOLD
bubble.line.width = Pt(2.5)
bubble.adjustments[0] = 0.1

add_text(s, "💬  Claude 대화창에 이렇게 적어주세요",
         1.7, 2.33, 10, 0.35, size=12, color=GRAY_MID)
add_text(s,
         "\"ironyjk/police-frameworks 플러그인을 저한테 설치해 주세요.\n"
         "경찰 프레임워크 12개가 들어있는 마켓플레이스입니다.\"",
         1.7, 2.7, 10, 1.1, size=17, bold=True, color=GRAY_DARK,
         align=PP_ALIGN.CENTER, line_spacing=1.45)

# 그 다음 설명
add_text(s, "그럼 Claude가 알아서",
         0.6, 4.3, 12, 0.4, size=16, bold=True, color=NAVY)

outcomes = [
    ("Claude Desktop 에서",
     "Claude가 Settings → Capabilities → Plugins 에서\n마켓플레이스를 추가하고 설치하는 방법을 화면 그대로 안내해 드립니다."),
    ("Claude 웹 채팅에서",
     "Customize → Skills 경로로 동일한 설치 절차를 안내합니다."),
    ("Claude Code (개발자용)에서",
     "`/plugin marketplace add ironyjk/police-frameworks` 명령어를\n자동 실행해 12개 공구를 바로 등록합니다."),
]

y = 4.8
for title, desc in outcomes:
    icon = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.75), Inches(y + 0.08), Inches(0.4), Inches(0.4)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = GREEN
    icon.line.fill.background()
    add_text(s, "✓", 0.75, y + 0.1, 0.4, 0.35,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.3, y + 0.08, 4, 0.4,
             size=13, bold=True, color=NAVY)
    add_text(s, desc, 5.4, y + 0.05, 7.5, 0.7,
             size=11, color=GRAY_DARK, line_spacing=1.3)
    y += 0.72

add_footer(s, 25)


# ══════════════════════════════════════════════════════════════
# Slide 26 — 첫 대화 예시
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, BG)
add_title_bar(s, "첫 대화 예시",
              "Claude 구독만 되어 있으면, 슬래시 명령어 몰라도 됩니다")

add_text(s, "Claude 대화창에 이렇게 적으시면 됩니다",
         0.6, 1.55, 12, 0.5, size=17, bold=True, color=NAVY)

chat_box = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(2.15), Inches(12.1), Inches(2.8)
)
chat_box.fill.solid()
chat_box.fill.fore_color.rgb = GRAY_PALE
chat_box.line.color.rgb = NAVY
chat_box.line.width = Pt(1.5)
chat_box.adjustments[0] = 0.05

add_text(s, "  💬  대화창 입력",
         0.75, 2.25, 11.8, 0.35, size=11, color=GRAY_MID)

example = (
    "peel 로 아래 민원을 분석해 주세요.\n\n"
    "경주 성건동 근린공원에서 3개월째 야간에 청소년 집단 음주·소란\n"
    "민원이 반복되고 있습니다. 14건 누적. 순찰 가면 흩어지고 돌아오면\n"
    "다시 모입니다. 한쪽 주민은 단속 강화를 요구하고, 다른 쪽은 낙인을\n"
    "반대합니다. 우리 지구대에서 어떻게 접근할까요?"
)
add_text(s, example, 0.9, 2.65, 11.8, 2.2,
         size=13, color=GRAY_DARK, line_spacing=1.45)

add_text(s, "그러면 /peel 이 이런 초안을 만들어 드립니다",
         0.6, 5.1, 12, 0.4, size=15, bold=True, color=NAVY)

bullets = [
    "선정 공구: SARA · Crime Triangle · Hot Spots · CPTED · COP · Procedural Justice · Broken Windows(비판)",
    "제외 공구: PEACE·Cognitive (조사 단계 아님) / BCSM·ICAT (위기 아님) / ILP (단일 지점)",
    "6가지 실행 방안 — 담당·기한·단계까지 구체적으로 (오늘 보신 사례 C 슬라이드와 동일)",
]
add_bullets(s, bullets, 0.85, 5.5, 12.1, 1.5, size=12)

add_footer(s, 26)


# ══════════════════════════════════════════════════════════════
# Slide 27 — 부탁과 연락처
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_background(s, NAVY_DEEP)

top_line = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.8), prs.slide_width, Inches(0.08)
)
top_line.fill.solid()
top_line.fill.fore_color.rgb = GOLD
top_line.line.fill.background()

add_text(s, "부탁드리는 한 가지",
         0.8, 1.1, 12, 0.6, size=20, color=GOLD)
add_text(s, "써 보시고 고쳐주십시오.",
         0.8, 1.75, 12, 1.0, size=46, bold=True, color=WHITE)
add_text(s,
         "이 도구는 완성품이 아닙니다. 초안입니다.\n"
         "여러분의 현장 경험으로 틀린 부분을 지적해 주시면,\n"
         "그게 다음 버전에 반영됩니다. 그게 가장 큰 기여입니다.",
         0.8, 3.1, 12, 1.6, size=17, color=GRAY_LIGHT, line_spacing=1.45)

contact_box = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.0), Inches(11.8), Inches(1.7)
)
contact_box.fill.solid()
contact_box.fill.fore_color.rgb = NAVY
contact_box.line.color.rgb = GOLD
contact_box.line.width = Pt(1.5)
contact_box.adjustments[0] = 0.1

add_text(s, "만든 사람", 1.1, 5.2, 5, 0.4, size=13, color=GOLD)
add_text(s, "최희철", 1.1, 5.5, 5, 0.6, size=26, bold=True, color=WHITE)
add_text(s, "경주경찰서 경찰발전협의회 회원",
         1.1, 6.1, 5.5, 0.4, size=13, color=GRAY_LIGHT)
add_text(s, "평소 AI 도구를 만드는 개발자",
         1.1, 6.4, 5.5, 0.4, size=13, color=GRAY_LIGHT)

add_text(s, "설치 링크", 7.5, 5.2, 5, 0.4, size=13, color=GOLD)
add_text(s, "ironyjk/police-", 7.5, 5.55, 5, 0.4,
         size=16, bold=True, color=WHITE)
add_text(s, "frameworks", 7.5, 5.85, 5, 0.4,
         size=16, bold=True, color=WHITE)
add_text(s, "Claude에게 이걸 깔아달라고만 하시면 됩니다.",
         7.5, 6.2, 5.8, 0.4, size=12, color=GRAY_LIGHT)

add_text(s,
         "\"경찰은 시민이고, 시민은 경찰이다.\"  — Robert Peel, 1829",
         0.8, 6.95, 12, 0.4, size=13, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────
# 저장
# ──────────────────────────────────────────────────────────────

out = HERE.parent / "docs" / "intro.pptx"
out.parent.mkdir(exist_ok=True)
prs.save(str(out))
print(f"OK Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
